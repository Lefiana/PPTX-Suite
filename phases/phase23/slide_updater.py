"""
phases/phase23/slide_updater.py
Business-logic layer for Phase 2+3.  Combines core.image_processor and
core.pptx_engine into the operations the QA controller actually needs:
roster loading, draft generation, single-slide re-injection, and preview
cropping.  No tkinter imports — fully unit-testable.
"""
from __future__ import annotations
from typing import Callable, Optional
import logging
from pathlib import Path
from PIL import Image
from pptx import Presentation

from core.image_processor import ImageProcessor
from core.pptx_engine import PPTXEngine

logger = logging.getLogger(__name__)


class SlideUpdater:
    """Stateless orchestration wrapper — construct once and reuse."""

    def __init__(self, config_manager) -> None:
        self.cfg = config_manager
        self.image_processor = ImageProcessor()
        self.pptx_engine     = PPTXEngine(self.image_processor, config_manager)

    # ── Roster ────────────────────────────────────────────────────────────────

    def load_students_from_excel(self, excel_path: str, master_dir: str) -> list[dict]:
        layout = self.cfg.load_layout_config()
        return self.pptx_engine.load_students_from_excel(excel_path, master_dir, layout)

    # ── Draft generation ──────────────────────────────────────────────────────

    def generate_draft(
        self,
        students:      list[dict],
        template_path: str,
        output_path:   str,
        on_progress:   Optional[Callable[[int, int], None]] = None,
    ) -> list[dict]:
        layout = self.cfg.load_layout_config()
        fd_cfg = layout.get("face_detection", {})
        return self.pptx_engine.generate_draft(
            students, template_path, output_path, layout, fd_cfg, on_progress
        )

    # ── Preview ───────────────────────────────────────────────────────────────

    def crop_preview(
        self,
        image_path: str,
        manual_face_center: Optional[tuple] = None,
        preview_size: tuple = (280, 380),
    ) -> Image.Image:
        layout = self.cfg.load_layout_config()
        ratio  = layout["portrait"]["width_cm"] / layout["portrait"]["height_cm"]
        fd_cfg = layout.get("face_detection", {})
        return self.image_processor.get_preview_pil(
            image_path, ratio, fd_cfg, manual_face_center, preview_size
        )

    def get_original_resized(self, image_path: str, max_size: tuple = (380, 500)):
        return self.image_processor.get_original_resized(image_path, max_size)

    def target_ratio(self) -> float:
        p = self.cfg.load_layout_config()["portrait"]
        return p["width_cm"] / p["height_cm"]

    # ── Re-injection ──────────────────────────────────────────────────────────

    def reinject_portrait(
        self,
        pptx_path:          str,
        slide_index:        int,
        portrait_shape_id:  int,
        image_path:         str,
        manual_face_center: Optional[tuple],
    ) -> int:
        """
        Crops image_path according to manual_face_center (or auto-detects if
        None), surgically replaces the portrait shape in pptx_path, and
        returns the new shape_id.
        """
        layout = self.cfg.load_layout_config()
        ratio  = layout["portrait"]["width_cm"] / layout["portrait"]["height_cm"]
        fd_cfg = layout.get("face_detection", {})

        stream = self.image_processor.crop_to_stream(
            image_path, ratio, fd_cfg, manual_face_center
        )
        return self.replace_slide_image(pptx_path, slide_index, portrait_shape_id, stream)

    def replace_slide_image(self, pptx_path, slide_index, portrait_shape_id, image_stream) -> int:
        path = Path(pptx_path)
        if not path.exists():
            logger.error("replace_slide_image: PPTX not found: %s", pptx_path)
            raise FileNotFoundError(f"PPTX not found: {pptx_path}")
        
        try:
            prs = Presentation(pptx_path)
        except Exception as exc:
            logger.exception("replace_slide_image: failed to open PPTX %s", pptx_path)
            raise RuntimeError(f"Could not open PPTX (it may be corrupted): {exc}") from exc

        if not (0 <= slide_index < len(prs.slides)):
            logger.error("replace_slide_image: slide_index %s out of range (deck has %s slides) in %s",
                         slide_index, len(prs.slides), pptx_path)
            raise IndexError(f"Slide index {slide_index} out of range (deck has {len(prs.slides)} slides).")

        slide = prs.slides[slide_index]
        if not any(shape.shape_id == portrait_shape_id for shape in slide.shapes):
            logger.error("replace_slide_image: shape_id %s not found on slide %s in %s",
                         portrait_shape_id, slide_index, pptx_path)
            raise ValueError(f"Portrait shape id {portrait_shape_id} not found on slide {slide_index}.")

        layout = self.cfg.load_layout_config()
        try:
            return self.pptx_engine.reinject_image(
                pptx_path, slide_index, portrait_shape_id, image_stream, layout)
        except Exception as exc:
            logger.exception("replace_slide_image: reinjection failed for slide %s in %s",
                             slide_index, pptx_path)
            raise RuntimeError(f"Portrait replacement failed: {exc}") from exc

    def update_student_slide(self, student: dict, image_stream) -> int:
        missing = [k for k in ("pptx_path", "slide_index", "portrait_shape_id") if k not in student]
        if missing:
            raise KeyError(f"Student record missing required field(s): {missing}")
        return self.replace_slide_image(
            student["pptx_path"], student["slide_index"],
            student["portrait_shape_id"], image_stream,
        )