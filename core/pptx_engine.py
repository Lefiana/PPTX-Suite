"""
core/pptx_engine.py
All python-pptx operations:
  • Slide duplication from a template
  • Text-frame replacement preserving original font style
  • Full draft-deck generation (thread-safe: returns updated student list)
  • Surgical per-slide image re-injection for live QA overrides
  • Excel → student-record parsing using layout_config schema
"""
from __future__ import annotations

import copy
import io
import re
import unicodedata
from pathlib import Path
from typing import Callable

import pandas as pd
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Cm

# ── Name utilities ────────────────────────────────────────────────────────────

def _normalize(text: str) -> str:
    t = unicodedata.normalize("NFKD", str(text)).encode("ASCII", "ignore").decode("ascii")
    return re.sub(r"[^\w\s]", " ", re.sub(r"\s+", " ", t)).lower().strip()


_SURNAME_PREFIXES = frozenset({
    "de", "del", "de la", "dela", "san", "santa", "sto", "sto.",
})


def parse_name(full_name: str) -> tuple[str, str]:
    """Returns (SURNAME, Firstname) from 'Lastname, Firstname' or plain 'Firstname Lastname'."""
    if "," in full_name:
        parts = full_name.split(",", 1)
        return parts[0].strip().upper(), parts[1].strip().title()

    words = full_name.split()
    if len(words) == 1:
        return words[0].upper(), ""

    if len(words) >= 2 and words[-2].lower() in _SURNAME_PREFIXES:
        return " ".join(words[-2:]).upper(), " ".join(words[:-2]).title()
    if len(words) >= 3 and " ".join(words[-3:-1]).lower() in _SURNAME_PREFIXES:
        return " ".join(words[-3:]).upper(), " ".join(words[:-3]).title()

    return words[-1].upper(), " ".join(words[:-1]).title()


# ── PPTXEngine ────────────────────────────────────────────────────────────────

class PPTXEngine:
    def __init__(self, image_processor, config_manager) -> None:
        self.ip  = image_processor
        self.cfg = config_manager

    # ── Excel → student records ───────────────────────────────────────────────

    def load_students_from_excel(
        self,
        excel_path: str,
        master_dir: str,
        layout_config: dict,
    ) -> list[dict]:
        exc      = layout_config.get("excel", {})
        header   = int(exc.get("header_row", 3))
        name_col = exc.get("name_column",    "STUDENT NAME")
        prog_col = exc.get("program_column", "PROGRAM")

        df     = pd.read_excel(excel_path, header=header)
        master = Path(master_dir) if master_dir else None

        col_map: dict[str, str] = {}
        for target in (name_col, prog_col):
            for col in df.columns:
                if str(col).strip().upper() == target.upper():
                    col_map[target] = col; break
            if target not in col_map:
                raise ValueError(
                    f"Column '{target}' not found.\nAvailable: {list(df.columns)}"
                )

        students: list[dict] = []
        for idx, row in df.iterrows():
            raw_name = str(row[col_map[name_col]]).strip()
            raw_prog = str(row[col_map[prog_col]]).strip()
            if not raw_name or raw_name.lower() in ("nan", "none", ""):
                continue
            if not raw_prog or raw_prog.lower() in ("nan", "none", ""):
                continue

            image_path = self._resolve_image(raw_name, raw_prog, master)
            surname, firstname = parse_name(raw_name)

            students.append({
                "excel_index":        idx,
                "surname":            surname,
                "firstname":          firstname,
                "course":             raw_prog,
                "image_path":         str(image_path),
                "slide_index":        -1,
                "portrait_shape_id":  -1,
                "status":             "placeholder" if "no_picture" in str(image_path) else "pending",
                "manual_face_center": None,
            })

        return sorted(students, key=lambda s: s["excel_index"])

    def _resolve_image(self, excel_name: str, program: str, master: Path | None) -> Path:
        fallback = Path("no_picture.png")
        if master is None or not master.exists():
            return fallback

        surname_part = excel_name.split(",")[0].strip() if "," in excel_name \
                       else excel_name.split()[-1].strip()
        norm_surname = _normalize(surname_part)
        norm_program = _normalize(program)

        program_dir: Path | None = None
        for p_dir in master.iterdir():
            if p_dir.is_dir() and norm_program in _normalize(p_dir.name):
                program_dir = p_dir; break

        if program_dir is None:
            return fallback

        for student_folder in program_dir.iterdir():
            if student_folder.is_dir() and norm_surname in _normalize(student_folder.name):
                for ext in ("*.jpg", "*.jpeg", "*.png"):
                    hit = next(student_folder.glob(ext), None)
                    if hit:
                        return hit
                break
        return fallback

    # ── Draft generation ──────────────────────────────────────────────────────

    def generate_draft(
        self,
        students:     list[dict],
        template_path: str,
        output_path:   str,
        layout_config: dict,
        fd_config:     dict,
        on_progress:   Callable[[int, int], None] | None = None,
    ) -> list[dict]:
        portrait = layout_config["portrait"]
        ratio    = portrait["width_cm"] / portrait["height_cm"]
        shapes   = layout_config.get("shapes", {})

        img_left   = Cm(portrait["left_cm"])
        img_top    = Cm(portrait["top_cm"])
        img_width  = Cm(portrait["width_cm"])
        img_height = Cm(portrait["height_cm"])

        prs   = Presentation(template_path)
        total = len(students)

        for count, student in enumerate(students):
            slide = self._duplicate_slide(prs, 0)

            for shape in slide.shapes:
                if shape.name == shapes.get("surname",   "TextBox 6"):
                    self._set_text(shape, student["surname"])
                elif shape.name == shapes.get("firstname", "TextBox 7"):
                    self._set_text(shape, student["firstname"])
                elif shape.name == shapes.get("course",    "TextBox 8"):
                    self._set_text(shape, student["course"])

            stream = self.ip.crop_to_stream(
                student["image_path"], ratio, fd_config,
                student.get("manual_face_center"),
            )
            pic = slide.shapes.add_picture(stream, img_left, img_top, img_width, img_height)

            student["slide_index"]      = count + 1
            student["portrait_shape_id"] = pic.shape_id
            student["status"] = (
                "placeholder" if "no_picture" in student["image_path"] else "generated"
            )

            if on_progress:
                on_progress(count + 1, total)

        # Remove template slide (index 0) — all generated slides shift down by 1
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]
        for student in students:
            student["slide_index"] -= 1

        prs.save(output_path)
        return students

    # ── Per-slide re-injection ─────────────────────────────────────────────────

    def reinject_image(
        self,
        pptx_path:         str,
        slide_index:       int,
        portrait_shape_id: int,
        image_stream:      io.BytesIO,
        layout_config:     dict,
    ) -> int:
        """Removes the old portrait shape and injects a new one. Returns the new shape_id."""
        portrait   = layout_config["portrait"]
        img_left   = Cm(portrait["left_cm"])
        img_top    = Cm(portrait["top_cm"])
        img_width  = Cm(portrait["width_cm"])
        img_height = Cm(portrait["height_cm"])

        prs   = Presentation(pptx_path)
        slide = prs.slides[slide_index]

        for shape in list(slide.shapes):
            if shape.shape_id == portrait_shape_id:
                shape.element.getparent().remove(shape.element)
                break

        pic = slide.shapes.add_picture(image_stream, img_left, img_top, img_width, img_height)
        prs.save(pptx_path)
        return pic.shape_id

    # ── PPTX helpers ──────────────────────────────────────────────────────────

    @staticmethod
    def _duplicate_slide(prs: Presentation, index: int):
        source    = prs.slides[index]
        new_slide = prs.slides.add_slide(source.slide_layout)
        for shp in list(new_slide.shapes):
            shp.element.getparent().remove(shp.element)
        for shape in source.shapes:
            new_slide.shapes._spTree.append(copy.deepcopy(shape.element))
        for rel in source.part.rels.values():
            if "notesSlide" not in rel.reltype:
                new_slide.part.relate_to(rel.target_part, rel.reltype)
        return new_slide

    @staticmethod
    def _set_text(shape, new_text: str) -> None:
        if not shape.has_text_frame:
            return
        shape.text_frame.word_wrap = False
        shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        first_para = shape.text_frame.paragraphs[0]
        run0 = first_para.runs[0] if first_para.runs else first_para.add_run()
        saved = {
            "name": run0.font.name, "size": run0.font.size,
            "bold": run0.font.bold, "italic": run0.font.italic,
        }
        try:
            saved["color"] = run0.font.color.rgb if run0.font.color.type else None
        except AttributeError:
            saved["color"] = None

        shape.text_frame.clear()
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = new_text
        run.font.name   = saved["name"]
        run.font.size   = saved["size"]
        run.font.bold   = saved["bold"]
        run.font.italic = saved["italic"]
        if saved["color"]:
            run.font.color.rgb = saved["color"]
